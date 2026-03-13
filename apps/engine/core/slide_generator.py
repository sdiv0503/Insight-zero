import io
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor  # <-- THE CORRECT IMPORT PATH

class BoardroomSlide:
    @staticmethod
    def generate_base64_slide(anomaly_date, revenue, confidence, root_cause) -> str:
        # 1. Create a blank presentation
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[5] # Title only layout
        slide = prs.slides.add_slide(blank_slide_layout)

        # 2. Add Title
        title_shape = slide.shapes.title
        title_shape.text = f"🚨 Insight-Zero Incident Report: {anomaly_date}"
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True

        # 3. Add Key Metrics Box
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        p = tf.add_paragraph()
        p.text = f"Impacted Revenue: ${revenue}"
        p.font.size = Pt(24)
        p.font.bold = True
        
        p2 = tf.add_paragraph()
        p2.text = f"AI Confidence Score: {confidence}"
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0, 102, 204) # Applies the blue color successfully

        # 4. Add AI Root Cause Box
        left = Inches(1)
        top = Inches(4)
        width = Inches(8)
        height = Inches(2.5)
        txBox2 = slide.shapes.add_textbox(left, top, width, height)
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        
        p3 = tf2.add_paragraph()
        p3.text = "Llama-3 Root Cause Analysis:"
        p3.font.size = Pt(20)
        p3.font.bold = True
        
        p4 = tf2.add_paragraph()
        p4.text = str(root_cause) if root_cause else "No AI analysis available."
        p4.font.size = Pt(16)

        # 5. Save to memory buffer instead of disk
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        # 6. Encode to Base64 so it easily travels over HTTP JSON
        slide_base64 = base64.b64encode(buffer.read()).decode("utf-8")
        return slide_base64